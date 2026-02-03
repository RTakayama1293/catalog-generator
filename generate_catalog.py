"""
商品カタログ生成スクリプト
テンプレートのプレースホルダーを置換して仕入先別カタログを生成

使い方:
    python generate_catalog.py HAK                    # 仕入先コード指定
    python generate_catalog.py HAK --excel data.xlsx  # Excelファイル指定
"""

import argparse
import copy
import io
import os
from pathlib import Path
from datetime import datetime

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


# === 設定 ===
DEFAULT_EXCEL = '受発注管理台帳.xlsx'
DEFAULT_TEMPLATE = 'catalog_template.pptx'
DEFAULT_IMAGES_DIR = 'images'
DEFAULT_OUTPUT_DIR = 'output'

PRODUCTS_PER_PAGE = 2


def extract_supplier_code(product_id):
    """商品連番から仕入先コードを抽出"""
    if pd.isna(product_id):
        return None
    parts = str(product_id).split('_')
    return parts[2] if len(parts) >= 3 else None


def safe_str(val, default='－'):
    """NaN安全な文字列変換"""
    if pd.isna(val):
        return default
    return str(val)


def format_price(val):
    """価格フォーマット"""
    if pd.isna(val):
        return '－'
    return f"¥{int(val):,}"


def load_data(excel_path, supplier_code):
    """Excelからデータ読み込み"""
    df = pd.read_excel(excel_path, sheet_name='商品マスタ', header=1)
    df['仕入先コード'] = df['商品連番'].apply(extract_supplier_code)
    
    products = df[(df['仕入先コード'] == supplier_code) & (df['商品名'].notna())].copy()
    supplier_name = products['仕入先'].iloc[0] if len(products) > 0 else supplier_code
    
    return products, supplier_name


def replace_text_in_paragraph(paragraph, replacements):
    """
    パラグラフ内のテキストを置換（run分割対応版）
    複数runに分割されたプレースホルダーも正しく置換する
    """
    # 全runのテキストを結合
    full_text = ''.join([run.text for run in paragraph.runs])
    
    # 置換が必要かチェック
    new_text = full_text
    for key, value in replacements.items():
        if key in new_text:
            new_text = new_text.replace(key, value)
    
    # 変更があった場合のみ書き戻す
    if new_text != full_text:
        # 最初のrunにテキスト全体を入れ、残りは空に
        if paragraph.runs:
            paragraph.runs[0].text = new_text
            for run in paragraph.runs[1:]:
                run.text = ''


def replace_text_in_shape(shape, replacements):
    """シェイプ内のテキストを置換"""
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            replace_text_in_paragraph(paragraph, replacements)


def replace_text_in_table(table, replacements):
    """テーブル内のテキストを置換"""
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                replace_text_in_paragraph(paragraph, replacements)


def convert_image_for_pptx(image_path):
    """画像をpython-pptx対応形式に変換（WebP→PNG）"""
    ext = os.path.splitext(image_path)[1].lower()
    if ext == '.webp':
        img = Image.open(image_path)
        img_bytes = io.BytesIO()
        img.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        return img_bytes
    return image_path


def find_and_replace_image(slide, placeholder_text, image_path):
    """画像プレースホルダーを実画像で置換（アスペクト比維持）"""
    if not os.path.exists(image_path):
        return False

    image_data = convert_image_for_pptx(image_path)

    # 画像のサイズを取得
    img = Image.open(image_path)
    img_width, img_height = img.size
    img_aspect = img_width / img_height

    for shape in slide.shapes:
        if shape.has_text_frame:
            # 全runを結合してチェック
            text = ''.join([
                ''.join([run.text for run in p.runs])
                for p in shape.text_frame.paragraphs
            ])
            if placeholder_text in text:
                placeholder_left, placeholder_top = shape.left, shape.top
                placeholder_width, placeholder_height = shape.width, shape.height
                placeholder_aspect = placeholder_width / placeholder_height

                # アスペクト比を維持してプレースホルダー内に収める
                if img_aspect > placeholder_aspect:
                    # 画像が横長 → 幅に合わせる
                    new_width = placeholder_width
                    new_height = int(placeholder_width / img_aspect)
                else:
                    # 画像が縦長 → 高さに合わせる
                    new_height = placeholder_height
                    new_width = int(placeholder_height * img_aspect)

                # 中央揃えのオフセット計算
                left = placeholder_left + (placeholder_width - new_width) // 2
                top = placeholder_top + (placeholder_height - new_height) // 2

                sp = shape._element
                sp.getparent().remove(sp)

                slide.shapes.add_picture(image_data, left, top, new_width, new_height)
                return True
    return False


def replace_image_placeholder_with_text(slide, placeholder_text, replacement_text):
    """画像プレースホルダーをテキストで置換（画像がない場合用）"""
    for shape in slide.shapes:
        if shape.has_text_frame:
            # 全runを結合してチェック
            text = ''.join([
                ''.join([run.text for run in p.runs])
                for p in shape.text_frame.paragraphs
            ])
            if placeholder_text in text:
                # プレースホルダーを置換テキストに変更
                for paragraph in shape.text_frame.paragraphs:
                    full_text = ''.join([run.text for run in paragraph.runs])
                    if placeholder_text in full_text:
                        new_text = full_text.replace(placeholder_text, replacement_text)
                        if paragraph.runs:
                            paragraph.runs[0].text = new_text
                            for run in paragraph.runs[1:]:
                                run.text = ''
                return True
    return False


def build_replacements(product, num, supplier_name):
    """商品データから置換辞書を生成"""
    # 列名（改行を含む）
    price_col = '国内定価\n（15％）'
    msrp_col = '参考上代\n（税込)'

    msrp_val = product[msrp_col]
    msrp_str = f"{format_price(msrp_val)}（税込）" if pd.notna(msrp_val) else '－'

    return {
        '{{仕入先名}}': supplier_name,
        f'{{{{商品名_{num}}}}}': safe_str(product['商品名']),
        f'{{{{容量_{num}}}}}': safe_str(product['容量']),
        f'{{{{単位_{num}}}}}': safe_str(product['単位']),
        f'{{{{MOQ_{num}}}}}': safe_str(product['発注ロット']),
        f'{{{{温度帯_{num}}}}}': safe_str(product['温度帯']),
        f'{{{{賞味期限_{num}}}}}': safe_str(product['賞味期限']),
        f'{{{{価格_{num}}}}}': format_price(product[price_col]),
        f'{{{{参考上代_{num}}}}}': msrp_str,
        f'{{{{商品説明_{num}}}}}': safe_str(product['商品特徴'], ''),
    }


def duplicate_slide(prs, slide_index):
    """スライドを複製（スライドマスター・レイアウトを維持）"""
    source_slide = prs.slides[slide_index]
    slide_layout = source_slide.slide_layout

    new_slide = prs.slides.add_slide(slide_layout)

    # 元スライドの要素をコピー
    for shape in source_slide.shapes:
        el = shape.element
        new_el = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    return new_slide


def generate_catalog(supplier_code, excel_path, template_path, images_dir, output_dir):
    """カタログ生成メイン処理"""

    # データ読み込み
    products, supplier_name = load_data(excel_path, supplier_code)
    print(f"仕入先: {supplier_name}")
    print(f"商品数: {len(products)}件")

    if len(products) == 0:
        print("エラー: 対象商品がありません")
        return None

    # テンプレートを直接使用（スライドマスター・ロゴを維持）
    prs = Presentation(template_path)

    # ページごとに処理（最初のページはテンプレートスライドを使用）
    for page_idx in range(0, len(products), PRODUCTS_PER_PAGE):
        page_products = products.iloc[page_idx:page_idx + PRODUCTS_PER_PAGE]

        if page_idx == 0:
            # 最初のページはテンプレートスライドをそのまま使用
            slide = prs.slides[0]
        else:
            # 2ページ目以降はテンプレートスライドを複製
            slide = duplicate_slide(prs, 0)
        
        # 置換辞書を構築
        replacements = {'{{仕入先名}}': supplier_name}
        
        for idx, (_, product) in enumerate(page_products.iterrows()):
            num = idx + 1
            replacements.update(build_replacements(product, num, supplier_name))
        
        # 2商品目がない場合は空欄に
        if len(page_products) < 2:
            for key in ['商品名', '容量', '単位', 'MOQ', '温度帯', '賞味期限', '価格', '参考上代', '商品説明']:
                replacements[f'{{{{{key}_2}}}}'] = ''
            replacements['{{画像_2}}'] = ''
        
        # テキスト置換
        for shape in slide.shapes:
            if shape.has_table:
                replace_text_in_table(shape.table, replacements)
            else:
                replace_text_in_shape(shape, replacements)
        
        # 画像置換
        for idx, (_, product) in enumerate(page_products.iterrows()):
            num = idx + 1
            product_id = product['商品連番']
            image_path = None

            # 対応形式: jpg, png, webp
            for ext in ['jpg', 'png', 'webp']:
                candidate = os.path.join(images_dir, f"{product_id}.{ext}")
                if os.path.exists(candidate):
                    image_path = candidate
                    break

            if image_path:
                find_and_replace_image(slide, f'{{{{画像_{num}}}}}', image_path)
            else:
                # 画像がない場合は "no image" を表示
                replace_image_placeholder_with_text(slide, f'{{{{画像_{num}}}}}', 'no image')
    
    # 保存
    os.makedirs(output_dir, exist_ok=True)
    date_str = datetime.now().strftime('%Y%m%d')
    output_filename = f"カタログ_{supplier_name}_{date_str}.pptx"
    output_path = os.path.join(output_dir, output_filename)

    prs.save(output_path)
    print(f"\n生成完了: {output_path}")
    
    return output_path


def main():
    parser = argparse.ArgumentParser(description='商品カタログ生成')
    parser.add_argument('supplier_code', help='仕入先コード（例: HAK）')
    parser.add_argument('--excel', default=DEFAULT_EXCEL, help='Excelファイルパス')
    parser.add_argument('--template', default=DEFAULT_TEMPLATE, help='テンプレートファイルパス')
    parser.add_argument('--images', default=DEFAULT_IMAGES_DIR, help='画像ディレクトリ')
    parser.add_argument('--output', default=DEFAULT_OUTPUT_DIR, help='出力ディレクトリ')
    
    args = parser.parse_args()
    
    generate_catalog(
        args.supplier_code,
        args.excel,
        args.template,
        args.images,
        args.output
    )


if __name__ == '__main__':
    main()
